"""Rebuild Logical Architecture (C4 L2) slide as NATIVE PowerPoint shapes.

Slide 10in x 5.625in (16:9).
* Title          y = 0.05 - 0.36
* Swimlanes      x = 0.00 - 7.85   (80%)
* Commentary     x = 8.00 - 9.95   (~20%)
* Footer         y = 5.45 - 5.60

Arrows are strictly orthogonal (horizontal + vertical only, 90 deg turns).
Idempotent - clears slide 5 and rebuilds.
"""
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

PPTX = r"D:\WorkspaceAI\hack2future\Hack2furure\docs\hackathon_blueprint_filled.pptx"


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


C = {
    "orange":    rgb("#ea580c"), "orange_d": rgb("#9a3412"), "orange_lt": rgb("#fff7ed"), "orange_bd": rgb("#fed7aa"),
    "blue":      rgb("#0078d4"), "blue_d":   rgb("#0a2540"), "blue_lt":   rgb("#dbeafe"), "blue_bg":   rgb("#eff6ff"), "blue_bd": rgb("#bfdbfe"),
    "purple":    rgb("#742774"), "purple_d": rgb("#3d2c70"), "purple_lt": rgb("#ede4f5"), "purple_bg": rgb("#faf5ff"), "purple_bd": rgb("#e9d5ff"),
    "green":     rgb("#107c10"), "green_d":  rgb("#0e3a0e"), "green_lt":  rgb("#dcfce7"), "green_bg":  rgb("#f0fdf4"), "green_bd": rgb("#bbf7d0"),
    "gold":      rgb("#d99416"), "gold_d":   rgb("#92400e"), "gold_lt":   rgb("#f5b73a"), "gold_bg":   rgb("#fef3c7"), "gold_bd":  rgb("#fde68a"),
    "red":       rgb("#dc2626"), "red_d":    rgb("#7f1d1d"), "red_lt":    rgb("#fee2e2"),
    "amber":     rgb("#d97706"),
    "white":     rgb("#ffffff"), "black":    rgb("#0f172a"), "grey":     rgb("#5a6680"),
}


def I(v):
    return Inches(v)


# ============== SHAPE HELPERS ==============
def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.75, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=8, bold=False, color=C["black"],
             align="left", valign="middle", italic=False):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.margin_left = Emu(15000); tf.margin_right = Emu(15000)
    tf.margin_top = Emu(5000); tf.margin_bottom = Emu(5000)
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, fill, line, line_w, text, *,
             size=8, bold=False, font_color=C["black"], align="center", valign="middle"):
    add_rect(slide, x, y, w, h, fill, line, line_w)
    add_text(slide, x, y, w, h, text, size=size, bold=bold, color=font_color,
             align=align, valign=valign)


def add_numbered_node(slide, x, y, w, h, fill, line, line_w, number, label,
                      *, num_size=12, label_size=6, num_color=C["black"], label_color=C["black"]):
    add_rect(slide, x, y, w, h, fill, line, line_w)
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.margin_left = Emu(10000); tf.margin_right = Emu(10000)
    tf.margin_top = Emu(8000); tf.margin_bottom = Emu(8000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = str(number)
    r1.font.size = Pt(num_size); r1.font.bold = True; r1.font.color.rgb = num_color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(label_size); r2.font.bold = True; r2.font.color.rgb = label_color


# ============== ORTHOGONAL ARROW HELPERS ==============
def _set_dashed(conn):
    ln = conn.line._get_or_add_ln()
    prstDash = ln.find(qn('a:prstDash'))
    if prstDash is None:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')


def _add_arrow_head(conn):
    ln = conn.line._get_or_add_ln()
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle'); tailEnd.set('w', 'sm'); tailEnd.set('len', 'sm')


def _draw_segment(slide, x1, y1, x2, y2, color, weight, dashed, arrow):
    """One straight line segment from (x1,y1) to (x2,y2)."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dashed:
        _set_dashed(conn)
    if arrow:
        _add_arrow_head(conn)
    return conn


def ortho(slide, x1, y1, x2, y2, color, weight=1.25, dashed=False, route="auto"):
    """Draw an orthogonal arrow from (x1,y1) to (x2,y2).

    route options:
      "auto"     - straight line if aligned on one axis, else pick longer axis first
      "h"        - horizontal first, then vertical
      "v"        - vertical first, then horizontal
      "hvh"      - horizontal + vertical + horizontal (3 segments, via mid_x)
      "vhv"      - vertical + horizontal + vertical (3 segments, via mid_y)
    """
    TOL = 0.005   # treat <0.005in as aligned
    aligned_x = abs(x1 - x2) < TOL
    aligned_y = abs(y1 - y2) < TOL

    if aligned_x or aligned_y:
        _draw_segment(slide, x1, y1, x2, y2, color, weight, dashed, arrow=True)
        return

    if route == "auto":
        # choose so the arrow ENTERS the target on its short axis
        route = "v" if abs(y2 - y1) > abs(x2 - x1) else "h"

    if route == "h":
        # horizontal first to (x2,y1), then vertical to (x2,y2)
        _draw_segment(slide, x1, y1, x2, y1, color, weight, dashed, arrow=False)
        _draw_segment(slide, x2, y1, x2, y2, color, weight, dashed, arrow=True)
    elif route == "v":
        # vertical first to (x1,y2), then horizontal to (x2,y2)
        _draw_segment(slide, x1, y1, x1, y2, color, weight, dashed, arrow=False)
        _draw_segment(slide, x1, y2, x2, y2, color, weight, dashed, arrow=True)
    elif route == "hvh":
        mid_x = (x1 + x2) / 2
        _draw_segment(slide, x1, y1, mid_x, y1, color, weight, dashed, arrow=False)
        _draw_segment(slide, mid_x, y1, mid_x, y2, color, weight, dashed, arrow=False)
        _draw_segment(slide, mid_x, y2, x2, y2, color, weight, dashed, arrow=True)
    elif route == "vhv":
        mid_y = (y1 + y2) / 2
        _draw_segment(slide, x1, y1, x1, mid_y, color, weight, dashed, arrow=False)
        _draw_segment(slide, x1, mid_y, x2, mid_y, color, weight, dashed, arrow=False)
        _draw_segment(slide, x2, mid_y, x2, y2, color, weight, dashed, arrow=True)


# ============== NODE REGISTRY ==============
NODES = {}

def register(key, x, y, w, h):
    NODES[key] = dict(x=x, y=y, w=w, h=h,
                      cx=x + w/2, cy=y + h/2,
                      top=(x + w/2, y),
                      bottom=(x + w/2, y + h),
                      left=(x, y + h/2),
                      right=(x + w, y + h/2))

def E(key, side, offset=0):
    """Return (x, y) of an edge point on a registered node, with optional offset."""
    n = NODES[key]
    if side == "top":    return (n["x"] + n["w"]/2 + offset, n["y"])
    if side == "bottom": return (n["x"] + n["w"]/2 + offset, n["y"] + n["h"])
    if side == "left":   return (n["x"], n["y"] + n["h"]/2 + offset)
    if side == "right":  return (n["x"] + n["w"], n["y"] + n["h"]/2 + offset)
    raise ValueError(side)


def slide_title_text(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.split("\n")[0].strip()
    return ""


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)


# ======================================================================
def build(slide):
    clear_slide(slide)
    NODES.clear()

    # ====== TITLE ======
    add_text(slide, 0.08, 0.04, 9.85, 0.20,
             "Logical Architecture (C4 L2)  —  swimlanes + numbered flow",
             size=12, bold=True, color=C["blue_d"], align="left", valign="middle")
    add_text(slide, 0.08, 0.22, 9.85, 0.14,
             "Numbers 1-13 trace runtime flow  ·  read with the commentary panel on the right",
             size=6, italic=True, color=C["grey"], align="left", valign="middle")

    # ====== SWIMLANES (left section, x=0 - 7.85) ======
    DIAG_RIGHT = 7.85
    LABEL_X = 0.04
    LABEL_W = 0.55
    CONT_X = 0.61
    CONT_W = DIAG_RIGHT - CONT_X - 0.02

    lane_specs = [
        ("CUSTOMER\n& CHANNELS",  0.40, 0.52, C["orange"], C["white"], C["orange_lt"], C["orange_bd"]),
        ("FRONTEND\n(SWA)",       0.94, 0.52, C["blue"],   C["white"], C["blue_bg"],   C["blue_bd"]),
        ("COPILOT\nSTUDIO",       1.48, 0.62, C["purple"], C["white"], C["purple_bg"], C["purple_bd"]),
        ("POWER\nAUTOMATE",       2.12, 0.62, C["purple"], C["white"], C["purple_bg"], C["purple_bd"]),
        ("AZURE AI\nSERVICES",    2.76, 0.92, C["blue"],   C["white"], C["blue_bg"],   C["blue_bd"]),
        ("M365 TEAMS\n(Human)",   3.70, 0.52, C["green"],  C["white"], C["green_bg"],  C["green_bd"]),
        ("DATAVERSE\nGlass Box",  4.24, 1.20, C["gold"],   C["white"], C["gold_bg"],   C["gold_bd"]),
    ]
    for label, y, h, lab_fill, lab_fc, bg_fill, bg_bd in lane_specs:
        add_rect(slide, CONT_X, y, CONT_W, h, bg_fill, bg_bd, 0.4, rounded=False)
        add_card(slide, LABEL_X, y, LABEL_W, h, lab_fill, lab_fill, 0.4,
                 label, size=6, bold=True, font_color=lab_fc, align="center", valign="middle")

    # ====== NUMBERED NODES — positions chosen so most arrows are pure H/V ======
    # Lane 1
    register("n1", 0.85, 0.44, 1.40, 0.42)
    add_numbered_node(slide, 0.85, 0.44, 1.40, 0.42, C["orange"], C["orange_d"], 0.5,
                      "1", "Customer + 5 channels",
                      num_color=C["white"], label_color=C["white"])

    # Lane 2  — n2 centered under n1, n3 to its right
    register("n2", 0.85, 0.98, 1.40, 0.42)
    add_numbered_node(slide, 0.85, 0.98, 1.40, 0.42, C["blue_lt"], C["blue"], 0.5,
                      "2", "React SPA + Token Broker",
                      num_color=C["blue_d"], label_color=C["blue_d"])
    register("n3", 2.55, 0.98, 1.30, 0.42)
    add_numbered_node(slide, 2.55, 0.98, 1.30, 0.42, C["blue_lt"], C["blue"], 0.5,
                      "3", "Direct Line / Bot",
                      num_color=C["blue_d"], label_color=C["blue_d"])

    # Lane 3 — n4 below n3 (same x col), n12 far right
    register("n4", 2.55, 1.55, 1.30, 0.48)
    add_numbered_node(slide, 2.55, 1.55, 1.30, 0.48, C["purple_lt"], C["purple"], 0.75,
                      "4", "Intake Agent (AI)",
                      num_color=C["purple_d"], label_color=C["purple_d"])
    register("n12", 6.20, 1.55, 1.55, 0.48)
    add_numbered_node(slide, 6.20, 1.55, 1.55, 0.48, C["purple_lt"], C["purple"], 0.75,
                      "12", "Explanation Agent (AI — NOT CSR)",
                      num_color=C["purple_d"], label_color=C["purple_d"])

    # Lane 4 — Create_Claim under n4, Master to its right, Router right, Notify under n12
    register("n5", 2.55, 2.20, 1.05, 0.48)
    add_numbered_node(slide, 2.55, 2.20, 1.05, 0.48, C["purple_lt"], C["purple"], 0.5,
                      "5", "Create_Claim",
                      num_color=C["purple_d"], label_color=C["purple_d"])
    register("n6", 3.75, 2.20, 1.25, 0.48)
    add_numbered_node(slide, 3.75, 2.20, 1.25, 0.48, C["purple_lt"], C["purple"], 0.75,
                      "6", "Master Orchestration",
                      num_color=C["purple_d"], label_color=C["purple_d"])
    # Tier router diamond
    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, I(5.15), I(2.20), I(0.85), I(0.48))
    diamond.fill.solid(); diamond.fill.fore_color.rgb = C["white"]
    diamond.line.color.rgb = C["green"]; diamond.line.width = Pt(1.0)
    diamond.shadow.inherit = False
    register("n11", 5.15, 2.20, 0.85, 0.48)
    add_text(slide, 5.15, 2.22, 0.85, 0.44, "11\nRouter",
             size=7, bold=True, color=C["green_d"], align="center", valign="middle")
    register("n13", 6.20, 2.20, 1.55, 0.48)
    add_numbered_node(slide, 6.20, 2.20, 1.55, 0.48, C["green_lt"], C["green"], 0.75,
                      "13", "Notify_Customer",
                      num_color=C["green_d"], label_color=C["green_d"])

    # Lane 5 — 3 agents stacked + Adjudication right
    register("n7", 2.55, 2.82, 1.75, 0.24)
    add_card(slide, 2.55, 2.82, 1.75, 0.24, C["blue_lt"], C["blue"], 0.5,
             "7  Extraction", size=7, bold=True, font_color=C["blue_d"], align="center", valign="middle")
    register("n8", 2.55, 3.08, 1.75, 0.24)
    add_card(slide, 2.55, 3.08, 1.75, 0.24, C["blue_lt"], C["blue"], 0.5,
             "8  Policy", size=7, bold=True, font_color=C["blue_d"], align="center", valign="middle")
    register("n9", 2.55, 3.34, 1.75, 0.24)
    add_card(slide, 2.55, 3.34, 1.75, 0.24, C["blue_lt"], C["blue"], 0.5,
             "9  Validation · 7 ext checks", size=6, bold=True, font_color=C["blue_d"], align="center", valign="middle")
    register("n10", 4.50, 3.00, 1.50, 0.55)
    add_numbered_node(slide, 4.50, 3.00, 1.50, 0.55, C["blue_d"], C["blue_d"], 1.0,
                      "10", "Adjudication · GPT-4.1",
                      num_color=C["white"], label_color=C["white"])

    # Lane 6 — humans
    register("nT2", 4.50, 3.76, 1.50, 0.20)
    add_card(slide, 4.50, 3.76, 1.50, 0.20, C["gold_bg"], C["amber"], 0.5,
             "T2 · Adjuster Teams Card", size=6, bold=True, font_color=C["gold_d"], align="center", valign="middle")
    register("nT3", 6.20, 3.76, 1.55, 0.20)
    add_card(slide, 6.20, 3.76, 1.55, 0.20, C["red_lt"], C["red"], 0.5,
             "T3 · Live CSR Teams Chat", size=6, bold=True, font_color=C["red_d"], align="center", valign="middle")
    add_text(slide, 0.65, 3.78, 3.80, 0.16,
             "Tier 1 = AI auto-approve, no human", size=6, italic=True, color=C["grey"],
             align="left", valign="middle")

    # Lane 7 — Dataverse tables + audit
    tables = [
        ("Policy",            0.70, 0.65),
        ("Claim",             1.40, 0.65),
        ("Document",          2.10, 0.80),
        ("Communication",     2.95, 1.05),
        ("Adjuster",          4.05, 0.65),
        ("Vendor",            4.75, 0.65),
        ("ClaimVendorAssgmt", 5.45, 1.40),
    ]
    for name, x, w in tables:
        add_card(slide, x, 4.32, w, 0.24, C["purple"], C["purple_d"], 0.3,
                 name, size=5, bold=True, font_color=C["white"], align="center", valign="middle")
    register("audit", 0.65, 4.62, 7.15, 0.40)
    add_card(slide, 0.65, 4.62, 7.15, 0.40, C["gold_lt"], C["gold"], 1.5,
             "🛡  Decision_Rationale  ·  THE GLASS BOX  ·  every AI step writes 1 row  ·  4 · 7 · 8 · 9 · 10 · 12",
             size=7, bold=True, font_color=C["blue_d"], align="center", valign="middle")

    # ====== ARROWS — strictly orthogonal ======
    R = C["red"]; B = C["black"]; G = C["gold"]

    # (1)→(2)  - pure vertical (n1 and n2 share center x = 1.55)
    ortho(slide, *E("n1", "bottom"), *E("n2", "top"), R, 1.4)
    # (2)→(3)  - pure horizontal (same y=1.19)
    ortho(slide, *E("n2", "right"), *E("n3", "left"), R, 1.4)
    # (3)→(4)  - pure vertical (share center x = 3.20)
    ortho(slide, *E("n3", "bottom"), *E("n4", "top"), R, 1.4)
    # (4)→(5)  - pure vertical (share center x = 3.20 / 3.08 ≈ aligned via L)
    ortho(slide, *E("n4", "bottom"), *E("n5", "top"), R, 1.4, route="v")
    # (5)→(6)  - pure horizontal
    ortho(slide, *E("n5", "right"), *E("n6", "left"), R, 1.4)
    # (6) fan-out → (7) (8) (9). Master bottom is at (4.375, 2.68).
    # Route each via vertical-then-horizontal so they don't cross.
    ortho(slide, *E("n6", "bottom"), *E("n7", "left"), R, 1.0, route="v")
    ortho(slide, E("n6", "bottom")[0]-0.20, E("n6", "bottom")[1], *E("n8", "left"), R, 1.0, route="v")
    ortho(slide, E("n6", "bottom")[0]-0.40, E("n6", "bottom")[1], *E("n9", "left"), R, 1.0, route="v")
    # (7)(8)(9)→(10) join — orthogonal H then V
    ortho(slide, *E("n7", "right"), *E("n10", "left", -0.15), R, 1.0, route="h")
    ortho(slide, *E("n8", "right"), *E("n10", "left"),         R, 1.0, route="h")
    ortho(slide, *E("n9", "right"), *E("n10", "left",  0.15),  R, 1.0, route="h")
    # (10)→(11) — vertical up + slight horizontal
    ortho(slide, *E("n10", "top"), *E("n11", "bottom"), R, 1.4, route="v")
    # (11)→T2 (black, vertical down)
    ortho(slide, *E("n11", "bottom"), *E("nT2", "top"), B, 1.0, route="v")
    # (11)→T3 (black, route h then v)
    ortho(slide, *E("n11", "right"), *E("nT3", "top"), B, 1.0, route="v")
    # (11)→(12) — vertical up (red, T1 path)
    ortho(slide, *E("n11", "top"), *E("n12", "bottom"), R, 1.4, route="v")
    # (12)→(13) — vertical down (same x center column area)
    ortho(slide, *E("n12", "bottom"), *E("n13", "top"), R, 1.4, route="v")
    # (13)→customer loop back — long red dashed via top-right corner
    nx, ny = E("n13", "right")
    cx, cy = E("n1", "top")
    # Right, then up, then left to customer top
    ortho(slide, nx, ny, 7.95, ny, R, 1.0, dashed=True)   # right edge
    ortho(slide, 7.95, ny, 7.95, 0.20, R, 1.0, dashed=True)  # up to top band
    ortho(slide, 7.95, 0.20, cx, 0.20, R, 1.0, dashed=True)  # left across top
    ortho(slide, cx, 0.20, cx, cy, R, 1.0, dashed=True)      # down into n1 top

    # Audit arrows (gold dashed) - draw 3 representative arrows down to audit row
    ax, ay = E("n4", "left")
    ortho(slide, ax, ay, ax - 0.20, ay, G, 0.75, dashed=True)
    ortho(slide, ax - 0.20, ay, ax - 0.20, 4.62, G, 0.75, dashed=True)
    ax2, ay2 = E("n10", "bottom")
    ortho(slide, ax2, ay2, ax2, 4.62, G, 0.75, dashed=True)
    ax3, ay3 = E("n12", "bottom")
    ortho(slide, ax3, ay3, ax3, 4.10, G, 0.75, dashed=True, route="v")
    ortho(slide, ax3, 4.10, 7.40, 4.10, G, 0.75, dashed=True)
    ortho(slide, 7.40, 4.10, 7.40, 4.62, G, 0.75, dashed=True)

    # ====== COMMENTARY PANEL (right ~20%, x=8.0 to 9.95) ======
    PX = 8.00; PW = 1.95; PY = 0.40; PH = 5.00
    add_rect(slide, PX, PY, PW, PH, C["white"], C["blue_d"], 1.0)
    add_card(slide, PX, PY, PW, 0.24, C["blue_d"], C["blue_d"], 0.5,
             "FLOW LEGEND", size=8, bold=True, font_color=C["white"], align="center", valign="middle")

    items = [
        ("1",  "CUSTOMER FNOL",
         "Mobile · Web · Teams · SMS · Email", C["orange_lt"], C["orange_d"]),
        ("2",  "REACT SPA + TOKEN BROKER",
         "Static Web Apps · token mint", C["blue_bg"], C["blue"]),
        ("3",  "DIRECT LINE BRIDGE",
         "Bot Service ↔ Copilot Studio", C["blue_bg"], C["blue"]),
        ("4",  "INTAKE AGENT (AI)",
         "11 loss topics · sub-flows", C["purple_bg"], C["purple"]),
        ("5",  "CREATE_CLAIM",
         "Inserts Claim + Document", C["purple_bg"], C["purple"]),
        ("6",  "MASTER ORCHESTRATION",
         "Trigger on Claim · fan-out", C["purple_bg"], C["purple"]),
        ("7",  "EXTRACTION AGENT",
         "Doc Intel + GPT-4o Vision", C["blue_bg"], C["blue"]),
        ("8",  "POLICY AGENT",
         "AI Search vector RAG", C["blue_bg"], C["blue"]),
        ("9",  "VALIDATION AGENT",
         "7 ext checks parallel", C["blue_bg"], C["blue"]),
        ("10", "ADJUDICATION (GPT-4.1)",
         "Decides + assigns tier", C["blue_bg"], C["blue"]),
        ("11", "TIER ROUTER",
         "T1 auto · T2 adj · T3 CSR", C["green_bg"], C["green"]),
        ("12", "EXPLANATION (AI)",
         "Plain-English rationale", C["purple_bg"], C["purple"]),
        ("13", "NOTIFY_CUSTOMER",
         "Reply on original channel", C["green_bg"], C["green"]),
    ]
    item_h = 0.30
    gap = 0.02
    item_y = PY + 0.27
    for num, title, body, fill, num_col in items:
        add_rect(slide, PX + 0.05, item_y, PW - 0.10, item_h, fill, num_col, 0.3)
        tb = slide.shapes.add_textbox(I(PX + 0.07), I(item_y), I(PW - 0.14), I(item_h))
        tf = tb.text_frame
        tf.margin_left = Emu(10000); tf.margin_right = Emu(10000)
        tf.margin_top = Emu(8000); tf.margin_bottom = Emu(8000)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = f"{num}  "
        r1.font.size = Pt(8); r1.font.bold = True; r1.font.color.rgb = num_col
        r2 = p.add_run(); r2.text = title
        r2.font.size = Pt(5); r2.font.bold = True; r2.font.color.rgb = C["black"]
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r3 = p2.add_run(); r3.text = body
        r3.font.size = Pt(5); r3.font.color.rgb = C["black"]
        item_y += item_h + gap

    # Glass Box callout at the bottom of panel
    add_rect(slide, PX + 0.05, item_y + 0.04, PW - 0.10, 0.40, C["gold_bg"], C["gold"], 1.0)
    tb = slide.shapes.add_textbox(I(PX + 0.07), I(item_y + 0.04), I(PW - 0.14), I(0.40))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(10000); tf.margin_right = Emu(10000)
    tf.margin_top = Emu(8000); tf.margin_bottom = Emu(8000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "🛡 GLASS BOX"
    r.font.size = Pt(7); r.font.bold = True; r.font.color.rgb = C["gold_d"]
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r = p2.add_run()
    r.text = "Every AI step → 1 row in Decision_Rationale"
    r.font.size = Pt(5); r.font.color.rgb = C["gold_d"]

    # ====== FOOTER ======
    add_text(slide, 0.08, 5.45, 9.85, 0.15,
             "ARROWS:  red = runtime flow  ·  red dashed = loop back  ·  black = Tier-2/3 human branches  ·  gold dashed = audit write",
             size=6, italic=True, color=C["grey"], align="left", valign="middle")


def main():
    if not os.path.exists(PPTX):
        raise SystemExit(f"PPTX not found: {PPTX}")
    prs = Presentation(PPTX)
    target = None
    for s in prs.slides:
        if slide_title_text(s).startswith("Logical Architecture"):
            target = s; break
    if target is None:
        raise SystemExit("Could not find 'Logical Architecture' slide")
    build(target)
    prs.save(PPTX)
    print(f"Done. Slide rebuilt with {len(target.shapes)} native shapes (orthogonal arrows + 20% commentary).")


if __name__ == "__main__":
    main()
