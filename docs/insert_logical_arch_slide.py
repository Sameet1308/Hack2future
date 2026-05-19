"""Insert (or REFRESH) the logical architecture PNG in hackathon_blueprint_filled.pptx.

Idempotent. If a slide titled "Logical Architecture (C4 L2)" already exists,
its picture is replaced in-place (no slide deletion - avoids python-pptx
slide-deletion bug that leaves orphan parts).

If not present, a new slide is inserted at TARGET_POS with the PNG fitted
to slide width and centered vertically.

Slide size: 10in x 5.625in (16:9)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os

PPTX = r"D:\WorkspaceAI\hack2future\Hack2furure\docs\hackathon_blueprint_filled.pptx"
PNG  = r"D:\WorkspaceAI\hack2future\Hack2furure\docs\02_logical_architecture.drawio.png"
SLIDE_TITLE = "Logical Architecture (C4 L2)"
TARGET_POS = 5   # 1-indexed; only used when inserting fresh


def slide_title_text(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.split("\n")[0].strip()
    return ""


def remove_pictures_from_slide(slide):
    """Strip every picture shape (shape_type 13) from the slide."""
    sp_tree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE
            sp_tree.remove(shape._element)
            removed += 1
    return removed


def fit_image_to_slide(prs, png_path):
    """Return (x, y, w, h) in EMU that fits the PNG into the slide, centered."""
    sw, sh = prs.slide_width, prs.slide_height
    img = Image.open(png_path)
    pw, ph = img.size
    aspect = pw / ph
    # Try fit-to-width first
    fw = sw
    fh = int(sw / aspect)
    if fh > sh:
        fh = sh
        fw = int(sh * aspect)
    x = (sw - fw) // 2
    y = (sh - fh) // 2
    return x, y, fw, fh


def add_title_textbox(slide):
    tb = slide.shapes.add_textbox(Inches(0.12), Inches(0.04), Inches(6), Inches(0.32))
    tf = tb.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.text = SLIDE_TITLE
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(11)
    run.font.bold = True


def main():
    if not os.path.exists(PPTX):
        raise SystemExit(f"PPTX not found: {PPTX}")
    if not os.path.exists(PNG):
        raise SystemExit(f"PNG not found: {PNG}")

    prs = Presentation(PPTX)
    x, y, fw, fh = fit_image_to_slide(prs, PNG)

    # Look for existing slide
    existing = [s for s in prs.slides if slide_title_text(s) == SLIDE_TITLE]

    if existing:
        target = existing[0]
        n_pics = remove_pictures_from_slide(target)
        target.shapes.add_picture(PNG, x, y, width=fw, height=fh)
        print(f"REFRESH: replaced {n_pics} picture(s) on existing slide '{SLIDE_TITLE}'.")
    else:
        # Pick blank layout
        blank_layout = None
        for layout in prs.slide_layouts:
            if len(layout.placeholders) == 0:
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(blank_layout)
        new_slide.shapes.add_picture(PNG, x, y, width=fw, height=fh)
        add_title_textbox(new_slide)
        # Reorder: move last slide to TARGET_POS
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)
        last = slides_list[-1]
        xml_slides.remove(last)
        xml_slides.insert(TARGET_POS - 1, last)
        print(f"INSERT: added new slide '{SLIDE_TITLE}' at position {TARGET_POS}.")

    prs.save(PPTX)
    print(f"Saved. Image fit: {fw/914400:.2f}in x {fh/914400:.2f}in  at offset ({x/914400:.3f}in, {y/914400:.3f}in)")


if __name__ == "__main__":
    main()
