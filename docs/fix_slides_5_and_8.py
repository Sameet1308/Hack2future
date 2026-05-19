"""Two fixes in one run:

1) SLIDE 5 - restore the LOGICAL ARCHITECTURE PNG (with real Microsoft icons).
   Removes the native-shape rebuild I did earlier and puts the PNG back
   edge-to-edge on a full slide. Real icons preserved.

2) SLIDE 8 - rewrite REFERENCES with the actual Confluence URLs (clickable
   hyperlinks) plus instructions for requesting read access.

Idempotent and safe - operates only on the two target slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PPTX = r"D:\WorkspaceAI\hack2future\Hack2furure\docs\hackathon_blueprint_filled.pptx"
PNG  = r"D:\WorkspaceAI\hack2future\Hack2furure\docs\02_logical_architecture.drawio.png"


# ===== Confluence references =====
CONFLUENCE_BASE = "https://aieliteltm.atlassian.net/wiki/spaces/PM/pages"
REFS = [
    ("Product Requirements (PRD)  ·  R1–R25",
     f"{CONFLUENCE_BASE}/1802274"),
    ("Dataverse Schema  ·  8 tables, slim + JSON pattern",
     f"{CONFLUENCE_BASE}/1802251"),
    ("Architecture  ·  4 layers, channels → orchestration → agents → data",
     f"{CONFLUENCE_BASE}/1769476"),
    ("Architectural Decisions Log  ·  17+ ADRs",
     f"{CONFLUENCE_BASE}/2097154"),
    ("Intake Agent FNOL Spec  ·  11 loss types, sub-flows, escalation",
     f"{CONFLUENCE_BASE}/1638411"),
]


def slide_title_text(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.split("\n")[0].strip()
    return ""


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


C = {
    "navy":      rgb("#0a2540"),
    "blue":      rgb("#0078d4"),
    "grey":      rgb("#5a6680"),
    "white":     rgb("#ffffff"),
    "amber":     rgb("#d97706"),
    "amber_lt":  rgb("#fef3c7"),
    "amber_d":   rgb("#92400e"),
    "orange":    rgb("#ea580c"),
}


# ============== SLIDE 5 ==============
def fix_slide_5(prs):
    # Find Logical Architecture slide
    target = None
    for s in prs.slides:
        if slide_title_text(s).startswith("Logical Architecture"):
            target = s; break
    if target is None:
        raise SystemExit("Could not find Logical Architecture slide")

    if not os.path.exists(PNG):
        print(f"WARNING: PNG not found at {PNG} — slide 5 will be left untouched.")
        return

    clear_slide(target)

    # Add small title at top-left
    tb = target.shapes.add_textbox(Inches(0.12), Inches(0.05), Inches(7), Inches(0.32))
    tf = tb.text_frame
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.text = "Logical Architecture (C4 L2)"
    r = tf.paragraphs[0].runs[0]
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = C["navy"]

    # Fit image to slide width, center vertically
    sw, sh = prs.slide_width, prs.slide_height
    img = Image.open(PNG); pw, ph = img.size
    aspect = pw / ph
    fw = sw; fh = int(sw / aspect)
    if fh > sh:
        fh = sh; fw = int(sh * aspect)
    x = (sw - fw) // 2
    y = (sh - fh) // 2
    target.shapes.add_picture(PNG, x, y, width=fw, height=fh)
    print(f"Slide 5 restored - picture {fw/914400:.2f}in x {fh/914400:.2f}in (real Microsoft icons preserved).")


# ============== SLIDE 8 ==============
def add_hyperlink_run(paragraph, text, url, *, size=10, bold=False, color=rgb("#0563c1")):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    # underline
    r.font.underline = True
    # hyperlink
    r.hyperlink.address = url


def add_plain_run(paragraph, text, *, size=10, bold=False, color=C["navy"], italic=False):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color


def fix_slide_8(prs):
    # Find References slide
    target = None
    for s in prs.slides:
        if slide_title_text(s).lower().startswith("references"):
            target = s; break
    if target is None:
        raise SystemExit("Could not find References slide")

    clear_slide(target)

    # Title
    tb = target.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.45))
    tf = tb.text_frame
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "References & Supporting Material", size=22, bold=True, color=C["navy"])

    # Subtitle
    tb = target.shapes.add_textbox(Inches(0.4), Inches(0.68), Inches(9.2), Inches(0.25))
    tf = tb.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p,
                  "All artifacts live in two places — public GitHub repo (read-only) and Confluence (PM space, access on request).",
                  size=10, italic=True, color=C["grey"])

    # ===== LEFT COLUMN: GitHub + Diagrams + Schema =====
    LX = 0.4; LY = 1.10; LW = 4.55; LH = 4.10
    # Section header
    tb = target.shapes.add_textbox(Inches(LX), Inches(LY), Inches(LW), Inches(0.28))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "📁  GitHub Repository  (public, no access needed)",
                  size=12, bold=True, color=C["navy"])
    # Repo link
    tb = target.shapes.add_textbox(Inches(LX), Inches(LY + 0.30), Inches(LW), Inches(0.25))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_hyperlink_run(p, "github.com/Sameet1308/Hack2future",
                      "https://github.com/Sameet1308/Hack2future", size=11, bold=True)

    # Diagrams sub-section
    yy = LY + 0.65
    tb = target.shapes.add_textbox(Inches(LX), Inches(yy), Inches(LW), Inches(0.25))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "Diagrams  (open in app.diagrams.net):",
                  size=10, bold=True, color=C["navy"])
    diagrams = [
        ("02_logical_architecture.drawio  ·  C4 Level 2 swimlanes",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/diagrams/02_logical_architecture.drawio"),
        ("03_er_diagram.drawio  ·  8 entities + relationships",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/diagrams/03_er_diagram.drawio"),
        ("04_assignment_flows.drawio  ·  Handler + Vendor engines",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/diagrams/04_assignment_flows.drawio"),
        ("05_blueprint_flow.drawio  ·  half-slide solution flow",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/diagrams/05_blueprint_flow.drawio"),
        ("06_agent_architecture.drawio  ·  6 agents, parallel pattern",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/diagrams/06_agent_architecture.drawio"),
    ]
    yy += 0.25
    for label, url in diagrams:
        tb = target.shapes.add_textbox(Inches(LX + 0.15), Inches(yy), Inches(LW - 0.15), Inches(0.20))
        tf = tb.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        add_plain_run(p, "• ", size=9, color=C["grey"])
        add_hyperlink_run(p, label, url, size=9)
        yy += 0.22

    # Schema sub-section
    yy += 0.08
    tb = target.shapes.add_textbox(Inches(LX), Inches(yy), Inches(LW), Inches(0.25))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "Schema files:", size=10, bold=True, color=C["navy"])
    yy += 0.25
    for label, url in [
        ("glassbox_schema.csv  ·  logical + physical column names",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/schema/glassbox_schema.csv"),
        ("glassbox_schema.xlsx  ·  10-sheet color-coded workbook",
         "https://github.com/Sameet1308/Hack2future/blob/main/docs/schema/glassbox_schema.xlsx"),
    ]:
        tb = target.shapes.add_textbox(Inches(LX + 0.15), Inches(yy), Inches(LW - 0.15), Inches(0.20))
        tf = tb.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        add_plain_run(p, "• ", size=9, color=C["grey"])
        add_hyperlink_run(p, label, url, size=9)
        yy += 0.22

    # ===== RIGHT COLUMN: Confluence pages =====
    RX = 5.10; RY = 1.10; RW = 4.55
    # Section header
    tb = target.shapes.add_textbox(Inches(RX), Inches(RY), Inches(RW), Inches(0.28))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "📚  Confluence  (PM space  ·  access on request)",
                  size=12, bold=True, color=C["navy"])

    # Workspace URL
    tb = target.shapes.add_textbox(Inches(RX), Inches(RY + 0.30), Inches(RW), Inches(0.25))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_hyperlink_run(p, "aieliteltm.atlassian.net/wiki/spaces/PM",
                      "https://aieliteltm.atlassian.net/wiki/spaces/PM/overview",
                      size=11, bold=True)

    # 5 page links
    yy = RY + 0.65
    for title, url in REFS:
        tb = target.shapes.add_textbox(Inches(RX), Inches(yy), Inches(RW), Inches(0.42))
        tf = tb.text_frame; tf.margin_top = 0; tf.margin_bottom = 0; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        add_hyperlink_run(p, title, url, size=10, bold=True)
        yy += 0.44

    # ===== HOW TO REQUEST ACCESS (bottom strip) =====
    BY = 5.30
    tb = target.shapes.add_textbox(Inches(0.4), Inches(BY), Inches(9.2), Inches(0.28))
    tf = tb.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "🔑  How to request Confluence access",
                  size=12, bold=True, color=C["amber_d"])

    tb = target.shapes.add_textbox(Inches(0.4), Inches(BY + 0.30), Inches(9.2), Inches(0.85))
    tf = tb.text_frame; tf.margin_top = 0; tf.word_wrap = True

    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "1.  ", size=10, bold=True, color=C["amber_d"])
    add_plain_run(p,
                  "Open ",
                  size=10)
    add_hyperlink_run(p, "aieliteltm.atlassian.net/wiki/spaces/PM",
                      "https://aieliteltm.atlassian.net/wiki/spaces/PM/overview", size=10)
    add_plain_run(p,
                  "  →  Atlassian will prompt you to sign in with a Microsoft / Google / Atlassian account.",
                  size=10)

    p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "2.  ", size=10, bold=True, color=C["amber_d"])
    add_plain_run(p,
                  'Click any restricted page  →  use the "Request access" link.  Approval lands in the space admin\'s inbox.',
                  size=10)

    p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
    add_plain_run(p, "3.  ", size=10, bold=True, color=C["amber_d"])
    add_plain_run(p,
                  "Faster path  —  email the team lead with your Atlassian-account email:  ",
                  size=10)
    add_hyperlink_run(p, "sameetdandawate@gmail.com",
                      "mailto:sameetdandawate@gmail.com?subject=Glass%20Box%20AI%20Confluence%20access%20request",
                      size=10, bold=True)
    add_plain_run(p, "  (typical turnaround under 1 business day).", size=10)

    print("Slide 8 rewritten with 5 Confluence page links, 5 GitHub diagram links, 2 schema links, + 3-step access instructions.")


def main():
    if not os.path.exists(PPTX):
        raise SystemExit(f"PPTX not found: {PPTX}")
    prs = Presentation(PPTX)
    fix_slide_5(prs)
    fix_slide_8(prs)
    prs.save(PPTX)
    print(f"Saved {PPTX}")


if __name__ == "__main__":
    main()
