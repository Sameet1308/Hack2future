"""Insert a native-pptx flow diagram into hackathon_blueprint_filled.pptx as a new slide
between slide 3 (Technical Architecture) and slide 4 (Team Profile).

No external dependencies (drawio CLI / cairo / docker). Pure python-pptx shapes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree
from pathlib import Path

SRC = Path("docs/hackathon_blueprint_filled.pptx")

# Color palette (Microsoft brand-ish)
COL_AZURE     = RGBColor(0x00, 0x78, 0xD4)
COL_AZURE_BG  = RGBColor(0xDB, 0xEA, 0xFE)
COL_PP        = RGBColor(0x74, 0x27, 0x74)
COL_PP_BG     = RGBColor(0xED, 0xE4, 0xF5)
COL_TEAMS     = RGBColor(0x5B, 0x41, 0xA1)
COL_TEAMS_BG  = RGBColor(0xED, 0xE4, 0xF5)
COL_M365      = RGBColor(0x10, 0x7C, 0x10)
COL_M365_BG   = RGBColor(0xDC, 0xFC, 0xE7)
COL_AUDIT     = RGBColor(0xD9, 0x94, 0x16)
COL_AUDIT_BG  = RGBColor(0xFE, 0xF3, 0xC7)
COL_TIER1     = RGBColor(0x16, 0xA3, 0x4A)
COL_TIER1_BG  = RGBColor(0xDC, 0xFC, 0xE7)
COL_TIER2     = RGBColor(0xD9, 0x77, 0x06)
COL_TIER2_BG  = RGBColor(0xFE, 0xF3, 0xC7)
COL_TIER3     = RGBColor(0xDC, 0x26, 0x26)
COL_TIER3_BG  = RGBColor(0xFE, 0xE2, 0xE2)
COL_DARK      = RGBColor(0x0A, 0x25, 0x40)
COL_GRAY      = RGBColor(0x5A, 0x66, 0x80)
COL_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

def add_box(slide, left, top, width, height, text, border_color, fill_color, font_color=None, bold=True, font_size=8):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.color.rgb = border_color
    sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = font_color if font_color else border_color
            r.font.name = "Segoe UI"
    return sh

def add_diamond(slide, left, top, width, height, text, border_color, fill_color, font_size=8):
    sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.color.rgb = border_color
    sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = border_color
            r.font.name = "Segoe UI"
    return sh

def add_line(slide, x1, y1, x2, y2, color=COL_GRAY, weight=1.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # End arrowhead
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn

def add_text(slide, left, top, width, height, text, size=10, bold=False, color=COL_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return tb

def move_slide_to_position(prs, slide_idx, new_idx):
    """Move slide from current index to new index using XML manipulation."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    item = slides[slide_idx]
    sldIdLst.remove(item)
    sldIdLst.insert(new_idx, item)

def main():
    prs = Presentation(SRC)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    # find a layout that's mostly blank (Title and Content or Blank)
    layout = None
    for lay in prs.slide_layouts:
        if "blank" in lay.name.lower():
            layout = lay
            break
    if layout is None:
        layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(layout)

    # Slide dims
    SW = prs.slide_width   # 9144000 EMU = 10"
    SH = prs.slide_height  # 5143500 EMU = 5.62"
    M = Inches(0.2)        # margin

    # Title
    add_text(slide, Inches(0.2), Inches(0.05), Inches(9.6), Inches(0.35),
             "Architecture Flow — top-to-bottom",
             size=16, bold=True, color=COL_DARK, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(0.2), Inches(0.32), Inches(9.6), Inches(0.18),
             "Customer → 5 channels → Copilot Studio → Power Automate → 4 parallel agents → audit → tier routing → notify",
             size=8, bold=False, color=COL_GRAY, align=PP_ALIGN.LEFT)

    # Layout: 11 rows top-to-bottom, ~0.42" each
    # Main column: x=0.4 to 8.0 (7.6 wide). Sidebar: x=8.1 to 9.85 (1.75 wide)
    MAIN_W = Inches(7.6)
    SIDE_X = Inches(8.1)
    SIDE_W = Inches(1.75)

    BX = Inches(0.4)  # main col x-start
    ROW_H = Inches(0.34)
    GAP   = Inches(0.06)
    y = Inches(0.55)

    # Helper for centering: x = BX + (MAIN_W - w)/2
    def center_x(w):
        return BX + (MAIN_W - w) // 2

    # Row 1: Customer
    cust_w = Inches(1.2)
    cust = add_box(slide, center_x(cust_w), y, cust_w, ROW_H,
                   "Customer", COL_DARK, COL_AZURE_BG, font_size=9)
    y_cust_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 2: 5 Channels
    chan_w = Inches(1.4); chan_gap = Inches(0.12)
    total_chan = 5*chan_w + 4*chan_gap
    chan_x_start = BX + (MAIN_W - total_chan)//2
    channels = []
    for i, (name, border, fill) in enumerate([
        ("Mobile App", COL_AZURE, COL_AZURE_BG),
        ("Web Chat",   COL_AZURE, COL_AZURE_BG),
        ("Microsoft Teams", COL_TEAMS, COL_TEAMS_BG),
        ("SMS",        COL_AZURE, COL_AZURE_BG),
        ("Email",      COL_M365,  COL_M365_BG),
    ]):
        x = chan_x_start + i*(chan_w + chan_gap)
        ch = add_box(slide, x, y, chan_w, ROW_H, name, border, fill, font_size=8)
        channels.append(ch)
    y_chan_top = y
    y_chan_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 3: Presentation (4 boxes: SWA, Token Broker Func, Direct Line, Azure Comm Services)
    pres_w = Inches(1.7); pres_gap = Inches(0.15)
    pres_total = 4*pres_w + 3*pres_gap
    pres_x_start = BX + (MAIN_W - pres_total)//2
    pres_boxes = []
    for i, name in enumerate([
        "Azure Static Web Apps",
        "Azure Function · Token Broker",
        "Bot Service · Direct Line",
        "Azure Comm Services (SMS)",
    ]):
        x = pres_x_start + i*(pres_w + pres_gap)
        b = add_box(slide, x, y, pres_w, ROW_H, name, COL_AZURE, COL_AZURE_BG, font_size=7)
        pres_boxes.append(b)
    y_pres_top = y
    y_pres_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 4: Copilot Studio Intake Agent (single, centered)
    cs_w = Inches(3.0)
    cs = add_box(slide, center_x(cs_w), y, cs_w, ROW_H,
                 "Copilot Studio · Intake Agent (FNOL_Start + 11 child topics)",
                 COL_PP, COL_PP_BG, font_size=8)
    y_cs_top = y; y_cs_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 5: Power Automate Master Orchestration
    pa_w = Inches(3.0)
    pa = add_box(slide, center_x(pa_w), y, pa_w, ROW_H,
                 "Power Automate · Master Orchestration",
                 COL_PP, COL_PP_BG, font_size=8)
    y_pa_top = y; y_pa_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 6: 4 Agents in parallel (Extraction, Policy, Validation, Adjudication)
    ag_w = Inches(1.75); ag_gap = Inches(0.1)
    ag_total = 4*ag_w + 3*ag_gap
    ag_x_start = BX + (MAIN_W - ag_total)//2
    agents = []
    for i, (name, border, fill) in enumerate([
        ("Extraction Agent\nDoc Intelligence + Vision", COL_AZURE, COL_AZURE_BG),
        ("Policy Agent\nAI Search RAG + AOAI",          COL_AZURE, COL_AZURE_BG),
        ("Validation Agent\nNOAA + NHTSA + 5 sandbox",  COL_AUDIT, COL_AUDIT_BG),
        ("Adjudication Agent\nAzure OpenAI GPT-4.1",    COL_AZURE, COL_AZURE_BG),
    ]):
        x = ag_x_start + i*(ag_w + ag_gap)
        h = Inches(0.55)
        b = add_box(slide, x, y, ag_w, h, name, border, fill, font_size=7)
        agents.append(b)
    y_ag_top = y; y_ag_bot = y + Inches(0.55)
    y = y_ag_bot + GAP

    # Row 7: Audit Log (wide, gold)
    aud_w = Inches(6.5)
    aud = add_box(slide, center_x(aud_w), y, aud_w, ROW_H,
                  "★ Decision_Rationale audit log — every AI decision in plain English",
                  COL_AUDIT, COL_AUDIT_BG, font_color=COL_AUDIT, font_size=8)
    y_aud_top = y; y_aud_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 8: Dataverse
    dv_w = Inches(3.5)
    dv = add_box(slide, center_x(dv_w), y, dv_w, ROW_H,
                 "Microsoft Dataverse · 8 tables",
                 COL_PP, COL_PP_BG, font_size=8)
    y_dv_top = y; y_dv_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 9: Tier router (diamond) + 3 tier outputs
    diamond_w = Inches(1.4); diamond_h = Inches(0.4)
    tier1_w = Inches(1.85); tier2_w = Inches(1.85); tier3_w = Inches(1.85); tier_gap = Inches(0.06)
    # Diamond on left, 3 tiers to right
    total_row = diamond_w + tier_gap + tier1_w + tier_gap + tier2_w + tier_gap + tier3_w
    row_x = BX + (MAIN_W - total_row)//2
    diamond = add_diamond(slide, row_x, y, diamond_w, diamond_h,
                          "Tier?", COL_AUDIT, COL_AUDIT_BG, font_size=8)
    t1_x = row_x + diamond_w + tier_gap
    t2_x = t1_x + tier1_w + tier_gap
    t3_x = t2_x + tier2_w + tier_gap
    tier1 = add_box(slide, t1_x, y, tier1_w, ROW_H, "Tier 1\nAuto-approve", COL_TIER1, COL_TIER1_BG, font_color=COL_TIER1, font_size=7)
    tier2 = add_box(slide, t2_x, y, tier2_w, ROW_H, "Tier 2\nTeams Adaptive Card", COL_TIER2, COL_TIER2_BG, font_color=COL_TIER2, font_size=7)
    tier3 = add_box(slide, t3_x, y, tier3_w, ROW_H, "Tier 3\nLive CSR escalation", COL_TIER3, COL_TIER3_BG, font_color=COL_TIER3, font_size=7)
    y_tier_top = y; y_tier_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 10: Explanation Agent + Notify_Customer (side by side)
    expl_w = Inches(2.4); notif_w = Inches(2.4); ex_gap = Inches(0.2)
    total_ex = expl_w + ex_gap + notif_w
    ex_x_start = BX + (MAIN_W - total_ex)//2
    expl = add_box(slide, ex_x_start, y, expl_w, ROW_H,
                   "Copilot Studio · Explanation Agent",
                   COL_PP, COL_PP_BG, font_size=8)
    notif = add_box(slide, ex_x_start + expl_w + ex_gap, y, notif_w, ROW_H,
                    "Power Automate · Notify_Customer",
                    COL_PP, COL_PP_BG, font_size=8)
    y_ex_top = y; y_ex_bot = y + ROW_H
    y += ROW_H + GAP

    # Row 11: Power BI (bottom)
    bi_w = Inches(2.5)
    bi = add_box(slide, center_x(bi_w), y, bi_w, ROW_H,
                 "Power BI · Operations Dashboard",
                 COL_PP, COL_PP_BG, font_size=8)
    y_bi_top = y; y_bi_bot = y + ROW_H

    # ============ SIDEBAR: Security & Compliance ============
    sx = SIDE_X
    sy = Inches(0.55)
    side_title = add_text(slide, sx, sy, SIDE_W, Inches(0.22),
                          "Security & Compliance",
                          size=10, bold=True, color=COL_DARK, align=PP_ALIGN.CENTER)
    sy += Inches(0.28)
    # Boxes
    for label, border, fill in [
        ("Microsoft Entra ID\nSSO + RBAC + roles",      COL_AZURE, COL_AZURE_BG),
        ("Azure Key Vault\nSecrets + keys",             COL_AZURE, COL_AZURE_BG),
        ("Application Insights\nTelemetry",             COL_AZURE, COL_AZURE_BG),
        ("Power Platform DLP\nConnector policies",      COL_PP,    COL_PP_BG),
        ("Encryption\nAt rest + in transit (TLS 1.2+)", COL_AZURE, COL_AZURE_BG),
    ]:
        add_box(slide, sx, sy, SIDE_W, Inches(0.45), label, border, fill, font_size=7)
        sy += Inches(0.5)

    # Compliance mapping
    add_box(slide, sx, sy, SIDE_W, Inches(0.85),
            "Compliance Mapping\nColorado SB21-169\nNAIC AI Bulletin\nNY DFS Circular 7\nCA AB 2930 (pending)",
            border_color=RGBColor(0x7c, 0x3a, 0xed),
            fill_color=RGBColor(0xed, 0xe9, 0xfe),
            font_color=RGBColor(0x5b, 0x21, 0xb6),
            font_size=7)

    # ============ ARROWS (no labels — plain connectors) ============
    # All arrows from y_chan_top: customer → channels
    cx_cust = (cust.left + cust.left + cust.width) // 2
    cy_cust = cust.top + cust.height
    for ch in channels:
        x_ch = (ch.left + ch.left + ch.width) // 2
        y_ch = ch.top
        add_line(slide, cx_cust, cy_cust, x_ch, y_ch)

    # Channels → presentation (each channel goes to nearest presentation box)
    # Simplified: channels 1,2 → SWA; 3 → Direct Line; 4 → ACS; 5 → ACS
    pres_map = [pres_boxes[0], pres_boxes[0], pres_boxes[2], pres_boxes[3], pres_boxes[3]]
    for ch, pres in zip(channels, pres_map):
        x_ch = (ch.left + ch.left + ch.width) // 2
        y_ch = ch.top + ch.height
        x_p = (pres.left + pres.left + pres.width) // 2
        y_p = pres.top
        add_line(slide, x_ch, y_ch, x_p, y_p)

    # Presentation → Copilot Studio (all 4 → CS)
    cs_top_x = (cs.left + cs.left + cs.width) // 2
    cs_top_y = cs.top
    for pres in pres_boxes:
        x_p = (pres.left + pres.left + pres.width) // 2
        y_p = pres.top + pres.height
        add_line(slide, x_p, y_p, cs_top_x, cs_top_y)

    # CS → PA Master
    add_line(slide, cs_top_x, cs.top + cs.height,
             (pa.left + pa.left + pa.width) // 2, pa.top, weight=1.5)

    # PA Master → 4 agents
    pa_bot_x = (pa.left + pa.left + pa.width) // 2
    pa_bot_y = pa.top + pa.height
    for ag in agents:
        x_a = (ag.left + ag.left + ag.width) // 2
        y_a = ag.top
        add_line(slide, pa_bot_x, pa_bot_y, x_a, y_a)

    # 4 agents → audit (gold thick)
    aud_top_x = (aud.left + aud.left + aud.width) // 2
    aud_top_y = aud.top
    for ag in agents:
        x_a = (ag.left + ag.left + ag.width) // 2
        y_a = ag.top + ag.height
        add_line(slide, x_a, y_a, x_a, aud_top_y, color=COL_AUDIT, weight=1.5)

    # Audit → Dataverse
    add_line(slide, aud_top_x, aud.top + aud.height,
             (dv.left + dv.left + dv.width) // 2, dv.top, weight=1.5)

    # Dataverse → diamond
    add_line(slide, (dv.left + dv.left + dv.width) // 2, dv.top + dv.height,
             (diamond.left + diamond.left + diamond.width) // 2, diamond.top, weight=1.5)

    # Diamond → 3 tiers
    dia_x = (diamond.left + diamond.left + diamond.width) // 2
    dia_y = diamond.top + diamond.height // 2
    add_line(slide, diamond.left + diamond.width, dia_y, tier1.left, tier1.top + tier1.height//2, color=COL_TIER1, weight=1.5)
    add_line(slide, diamond.left + diamond.width, dia_y, tier2.left, tier2.top + tier2.height//2, color=COL_TIER2, weight=1.5)
    add_line(slide, diamond.left + diamond.width, dia_y, tier3.left, tier3.top + tier3.height//2, color=COL_TIER3, weight=1.5)

    # Tier 1 + Tier 2 → Explanation (Tier 3 escalates to live human, no auto explanation)
    expl_top_x = (expl.left + expl.left + expl.width) // 2
    expl_top_y = expl.top
    add_line(slide, (tier1.left + tier1.left + tier1.width) // 2, tier1.top + tier1.height, expl_top_x, expl_top_y)
    add_line(slide, (tier2.left + tier2.left + tier2.width) // 2, tier2.top + tier2.height, expl_top_x, expl_top_y)

    # Explanation → Notify_Customer (horizontal)
    add_line(slide, expl.left + expl.width, expl.top + expl.height//2,
             notif.left, notif.top + notif.height//2)

    # Notify → Power BI (Power BI also reads Dataverse separately, but main flow ends here)
    add_line(slide, (notif.left + notif.left + notif.width) // 2, notif.top + notif.height,
             (bi.left + bi.left + bi.width) // 2, bi.top)

    # ============ Move slide to position 3 (after current slide 3) ============
    move_slide_to_position(prs, len(prs.slides) - 1, 3)

    prs.save(SRC)
    print(f"Inserted flow diagram as slide 4 of {SRC}")
    print(f"Total slides now: {len(prs.slides)}")

if __name__ == "__main__":
    main()
