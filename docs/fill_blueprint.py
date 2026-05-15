"""Fill in docs/hackathon_blueprint_template.pptx with Glass Box AI content.
Saves to docs/hackathon_blueprint_filled.pptx (preserves the original template).

Slides 1-5 get our content. Slides 6-7 (sample) are left intact.
Architecture diagram on slide 3 is intentionally left for later (per user request).
"""
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from pathlib import Path
from lxml import etree

SRC = Path("docs/hackathon_blueprint_template.pptx")
DST = Path("docs/hackathon_blueprint_filled.pptx")

# ============ CONTENT ============
SLIDE1 = {
    "Team Name":                  "Glass Box AI",
    "Team SPOC":                  "Sameet Dandawate",
    "Chosen Challenge Statement": "Business Challenge 2 — Insurance Agentic Claims Processing Solution. "
                                  "Build an agentic AI system for end-to-end claim processing including "
                                  "multi-channel intake, policy validation, missing-information retrieval, "
                                  "document understanding, adjudication, human-in-loop handling, and CSR support.",
    "Idea Name":                  "Glass Box AI — agentic claim processing on Microsoft Power Platform + Azure AI",
}

SLIDE2 = {
    "Project Proposal":
        "PROBLEM: US auto-insurance claims average ~14 days to settle. The industry "
        "loses $40B+/yr to fraud. New AI regulations (Colorado SB21-169, NAIC AI Bulletin "
        "adopted by 20+ states, NY DFS Circular 7) penalize black-box AI decisions.\n"
        "\n"
        "SOLUTION: Glass Box AI is a 5-agent swarm + 2 assignment engines + 6th explanation "
        "agent, all orchestrated by Power Automate, with EVERY AI decision logged in plain "
        "English to a Dataverse Decision_Rationale audit table. Customer touchpoint = 'Sara' "
        "persona across 5 channels (Mobile App, Web Chat, Teams, SMS, Email).\n"
        "\n"
        "TARGET USERS: US Personal Auto policyholders (consumer side); claim handlers and "
        "live CSRs at the carrier (operations side); compliance officers and state "
        "insurance regulators (audit side).\n"
        "\n"
        "BUSINESS / TECHNICAL VALUE: settlement in <2 minutes for simple claims (vs "
        "~14 days industry avg) · 50%+ auto-approval rate · 7 parallel fraud-detection "
        "checks per claim · 100% of decisions audit-logged · regulatory compliance "
        "by design (not bolted on).\n"
        "\n"
        "REAL-WORLD APPLICATION: any US Personal Auto carrier. Architecture is "
        "production-ready; 6 industry-data adapters (ISO ClaimSearch, NICB, CARFAX, "
        "DMV, KBB, telematics) use sandbox endpoints today and swap to live during "
        "the carrier's standard 60-90 day procurement cycle.\n"
        "\n"
        "FUTURE EXTENSIBILITY: full claim lifecycle (5 phases incl. Reopen) covers all "
        "11 Personal Auto loss types end-to-end. Roadmap: WhatsApp channel, Fitbit "
        "health-sensor corroboration for PIP/MedPay, expand to Property and Specialty "
        "product lines, real ML fraud model (vs current rules-based), Voice/IVR via "
        "Dynamics 365 Contact Center.",

    "Microsoft Technology / Services Touched":
        "AI / DATA / APP: Microsoft Copilot Studio (Intake + Explanation agents) · "
        "Azure OpenAI Service (GPT-4.1 — adjudication, sentiment, explanation, vision) · "
        "Azure AI Document Intelligence (prebuilt invoice/receipt/idDocument/layout) · "
        "Azure AI Search (vector RAG over policy PDFs) · Microsoft Dataverse (8 tables, "
        "Decision_Rationale = compliance artifact) · Power Automate (Master Orchestration "
        "+ all sub-agent flows + Handler/Vendor Assignment Engines) · Power BI Embedded "
        "(operations dashboard, Teams-embedded).\n"
        "\n"
        "CHANNELS / INTEGRATION: Microsoft Teams (Adaptive Cards for adjusters, live "
        "CSR handover) · Office 365 Outlook (email channel + notifications) · Azure "
        "Communication Services (SMS) · Bot Framework Direct Line (web chat embed).\n"
        "\n"
        "INFRASTRUCTURE: Azure Static Web Apps (React frontend hosting + Azure Function "
        "token broker) · Azure Blob Storage (policy PDFs + uploaded photos) · Azure Key "
        "Vault (secrets) · Microsoft Entra ID (handler SSO, RBAC, app registrations).\n"
        "\n"
        "DEVOPS / SECURITY: Power Platform Pipelines (solution Dev → Test → Prod) · "
        "GitHub Actions (frontend CI/CD into Static Web Apps) · Application Insights + "
        "Power Platform Analytics (telemetry) · Power Platform DLP policies (connector "
        "classification).",

    "Business Impact":
        "MARKET: US Personal Auto premium = ~$310B/yr. ~10M auto claims filed annually. "
        "Industry leader Lemonade has shown sub-3-second AI claim settlement is feasible.\n"
        "\n"
        "COST BENEFITS: 50%+ Tier-1 auto-approval = 50%+ reduction in claim-handler "
        "workload for routine claims. 7-check parallel validation catches more fraud "
        "earlier in the funnel.\n"
        "\n"
        "PERFORMANCE METRICS:\n"
        "• Time to decision (Tier 1 auto-approve): <2 minutes vs ~14 days industry avg\n"
        "• FNOL completion time: <90 seconds end-to-end\n"
        "• Validation sub-checks per claim: 7 (NOAA, NHTSA, ISO, NICB, DMV, Telematics, EstimateRule)\n"
        "• Channels live: 5 (Mobile App, Web Chat, Teams, SMS, Email)\n"
        "• Loss-type coverage: all 11 Personal Auto loss types end-to-end\n"
        "• Audit-trail coverage: 100% of AI decisions logged to Decision_Rationale\n"
        "\n"
        "SECURITY & COMPLIANCE: Microsoft Entra ID SSO + RBAC for handlers · Azure Key "
        "Vault for all secrets · Power Platform DLP policies enforce connector boundaries · "
        "Decision_Rationale table = audit artifact for Colorado SB21-169 + NAIC AI Bulletin "
        "(20+ states) + NY DFS Circular Letter No. 7 + CA AB 2930 (pending).",
}

SLIDE3 = {
    "Architecture Overview":
        "4-layer architecture on Microsoft Power Platform + Azure AI:\n"
        "(1) Channels — 5 customer touchpoints (Mobile App, Web Chat, Teams, SMS, Email).\n"
        "(2) Presentation — Azure Static Web Apps hosts the React SPA serving both customer "
        "and adjuster surfaces, with an Azure Function token broker bridging Direct Line.\n"
        "(3) Conversation + Orchestration — Copilot Studio (Intake + Explanation Agents); "
        "Power Automate Master_Orchestration flow fans out into 5 parallel agents (Extraction, "
        "Policy, Validation, Adjudication) plus Handler + Vendor Assignment Engines.\n"
        "(4) Data + Audit — Microsoft Dataverse with 8 tables; Decision_Rationale is the "
        "audit / compliance artifact. Power BI dashboard reads directly via native connector.\n"
        "\n"
        "FNOL HAPPY-PATH (Tier 1 auto-approve, ~90 seconds end-to-end):\n"
        "1. Customer files claim via Mobile App\n"
        "2. SPA → Direct Line → Copilot Studio Intake Agent\n"
        "3. Sara walks through universal questions + Collision-specific branch\n"
        "4. Power Automate Create_Claim writes Claim row to Dataverse\n"
        "5. Dataverse trigger fires Master_Orchestration\n"
        "6. Parallel: Extraction (Doc Intelligence) · Policy (AI Search RAG) · Validation (NOAA + NHTSA live + 5 sandbox)\n"
        "7. Every step writes a Decision_Rationale row\n"
        "8. Adjudication Agent (Azure OpenAI) synthesizes → confidence + recommendation\n"
        "9. Tier 1 → auto-approve, Notify_Customer flow → Explanation Agent on original channel\n"
        "10. Power BI dashboard updates in real time",

    "Core Components":
        "CONVERSATIONAL: Copilot Studio Intake Agent (FNOL_Start parent + 11 child topics, one per loss type, + 3 reusable sub-flows: OtherParty, Witness, InjuryTriage). Copilot Studio Explanation Agent (post-resolution).\n"
        "\n"
        "ORCHESTRATION: Power Automate flows — Create_Claim · Master_Orchestration · Notify_Customer · Log_To_Audit (reusable child flow called by every agent). Sub-flows: Extraction · Policy · Validation (× 7 sub-checks) · Adjudication · Send_TeamsCard · Process_TeamsResponse · Assign_Claim_To_Handler · Assign_Vendors_To_Claim.\n"
        "\n"
        "AI SERVICES: Azure OpenAI (GPT-4.1) for adjudication, sentiment analysis, explanation generation, GPT-4o Vision for damage photos. Azure AI Search for vector RAG over policy PDFs. Azure AI Document Intelligence for prebuilt invoice/receipt/idDocument/layout extraction.\n"
        "\n"
        "DATA: Dataverse — 8 tables (Policy, Claim, Document, Communication, Decision_Rationale, Adjuster, Vendor, ClaimVendorAssignment). Slim universal columns + JSON column for loss-type-specific data. Azure Blob Storage for policy PDFs and uploaded files. Azure AI Search index over policy PDFs.\n"
        "\n"
        "USER-FACING: React Customer App (mobile-styled phone-frame Web Chat with Sara avatar) + React Adjuster Console (queue, claim detail, Theater Mode live agent visualization), both hosted on Azure Static Web Apps. Microsoft Teams Adaptive Cards for Tier-2 adjuster review. Live CSR handover via Teams chat for Tier-3.\n"
        "\n"
        "EXTERNAL DATA ADAPTERS: 2 LIVE — NOAA Weather (api.weather.gov) + NHTSA Recalls (vpic.nhtsa.dot.gov). 6 SANDBOX with production-final interfaces — ISO ClaimSearch, NICB, CARFAX, state DMV, KBB/NADA, telematics (UBI). Real endpoints configured during carrier's standard 60-90 day procurement.\n"
        "\n"
        "ANALYTICS: Power BI Embedded operations dashboard — direct Dataverse query, no ETL — embedded as a tile inside the Adjuster Console Teams channel.",

    "Security & Compliance":
        "IDENTITY: Microsoft Entra ID SSO for handler routes (claimsHandler / supervisor / admin roles) via Azure Static Web Apps built-in auth. App registrations for SWA + Power Platform service principal + per-AAD-secured Azure resource.\n"
        "\n"
        "SECRETS: Azure Key Vault holds DIRECT_LINE_SECRET, AAD client credentials, Azure OpenAI keys, Azure Communication Services connection string. Power Platform connection references store per-service OAuth/key connections, environment-scoped (Dev / Test / Prod isolated).\n"
        "\n"
        "DATA PROTECTION: All data encrypted at rest (Dataverse + Azure Storage default) and in transit (TLS 1.2+). Role-based access at the Dataverse table + record level. Audit logging via Dataverse audit log + Application Insights + Power Platform Analytics.\n"
        "\n"
        "GOVERNANCE: Power Platform DLP (Data Loss Prevention) policies classify connectors as Business / Non-business / Blocked. HTTP connector explicitly allowed for sandbox-adapter calls. ALM via Power Platform Pipelines (solution export Dev → Test → Prod) + GitHub Actions (frontend auto-deploy on push to main).\n"
        "\n"
        "REGULATORY COMPLIANCE — the killer pitch lever:\n"
        "• Colorado SB21-169 — Algorithmic non-discrimination + AI governance + explainability for insurers. Decision_Rationale answers all three.\n"
        "• NAIC AI Model Bulletin (Dec 2023, adopted by 20+ states) — AI governance, documentation, third-party vendor controls, monitoring. Same artifact.\n"
        "• NY DFS Circular Letter No. 7 (2024) — AI/ML use documentation + consumer challenge ability. Same artifact.\n"
        "• CA AB 2930 (pending) — Algorithmic decision-making transparency. Same artifact.\n"
        "\n"
        "Single architectural artifact (Decision_Rationale audit log) = single compliance artifact answering all four regulations. No separate compliance tooling needed.",
}

# Team Profile (slide 4) — 6 NAME positions
TEAM = [
    ("Sameet Dandawate",      "Team SPOC · Frontend + Architecture lead"),
    ("Prasad Barsode",         "Policy data + RAG · Confluence + PRD"),
    ("Suraj",                  "Telematics adapter · IoT roadmap"),
    ("[Member 4 — TBD]",       "Copilot Studio lead · Intake Agent topics"),
    ("[Member 5 — TBD]",       "Power Automate lead · Orchestration + Glass Box logging"),
    ("[Member 6 — TBD]",       "Demo + Power BI dashboard · backup video"),
]

SLIDE5_REFS = (
    "REPOSITORY · github.com/Sameet1308/Hack2future\n"
    "\n"
    "CONFLUENCE PAGES (Product management space, aieliteltm.atlassian.net):\n"
    "• Glass Box AI — Product Requirements (R1–R25) — PM/1802274\n"
    "• Glass Box AI — Dataverse Schema (8 tables, naming conventions) — PM/1802251\n"
    "• Glass Box AI — Architecture (4 layers, all Microsoft services) — PM/1769476\n"
    "• Glass Box AI — Architectural Decisions Log (17+ ADRs) — PM/2097154\n"
    "• Glass Box AI — Intake Agent FNOL Spec (11 loss types) — PM/1638411\n"
    "\n"
    "DIAGRAMS (editable in app.diagrams.net):\n"
    "• docs/diagrams/02_logical_architecture.drawio — C4 Level 2 with real Microsoft icons\n"
    "• docs/diagrams/03_er_diagram.drawio — Entity-Relationship diagram, 8 entities\n"
    "• docs/diagrams/04_assignment_flows.drawio — Handler + Vendor assignment engines\n"
    "• docs/diagrams/architecture.html — interactive end-to-end flow (browser-renderable)\n"
    "\n"
    "SCHEMA (editable formats):\n"
    "• docs/schema/glassbox_schema.csv — flat, importable, both logical + physical names\n"
    "• docs/schema/glassbox_schema.xlsx — 10-sheet color-coded workbook\n"
    "\n"
    "DEMO ASSETS:\n"
    "• Live React app (Vite + Tailwind) — frontend/ folder · npm install && npm run dev\n"
    "• Theater Mode (live agent visualization) — /handler/theater/CLM-2026-4520\n"
    "\n"
    "DELIVERY TIMELINE:\n"
    "• Blueprint submission: 2026-05-18\n"
    "• Blueprint evaluation: 2026-05-25 → 2026-05-29\n"
    "• Solution build window: 2026-06-01 → 2026-06-12 (12 days)\n"
    "• Solution judging: 2026-06-16 → 2026-06-24\n"
    "• Closure + winners announced: 2026-06-30 → 2026-07-07\n"
    "\n"
    "REGULATORY REFERENCES:\n"
    "• Colorado SB21-169 — AI/ML insurer non-discrimination + governance\n"
    "• NAIC AI Model Bulletin (Dec 2023) — adopted by 20+ states\n"
    "• NY DFS Circular Letter No. 7 (2024) — AI/ML documentation + consumer challenge\n"
    "• CA AB 2930 (pending) — algorithmic decision-making transparency\n"
    "\n"
    "REFERENCE REPOSITORIES (Microsoft + open-source):\n"
    "• microsoft/claims-processing-hack — Microsoft's own multi-agent claims hackathon\n"
    "• MSUSAzureAccelerators/AI-Powered-Insurance-Claims-Automation-Accelerator\n"
    "• databricks-industry-solutions/smart-claims (cross-stack patterns)"
)

# ============ HELPERS ============

def set_cell_text(cell, text):
    """Replace a table cell's text while preserving the first paragraph's formatting."""
    tf = cell.text_frame
    # Keep the first run's formatting if present
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        bold = first_run.font.bold
    else:
        font_name = font_size = bold = None
    # Clear existing paragraphs by removing them (keep one)
    txBody = tf._txBody
    # Remove all <a:p> children except the first
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    paras = txBody.findall('a:p', nsmap)
    for p in paras[1:]:
        txBody.remove(p)
    first_p = paras[0] if paras else None
    if first_p is not None:
        # Remove all runs from first_p
        for r in first_p.findall('a:r', nsmap):
            first_p.remove(r)
        for br in first_p.findall('a:br', nsmap):
            first_p.remove(br)

    # Now write new text, splitting on \n into multiple paragraphs
    lines = text.split("\n")
    # First paragraph: replace
    p0 = tf.paragraphs[0]
    run = p0.add_run()
    run.text = lines[0]
    if font_name: run.font.name = font_name
    if font_size: run.font.size = font_size
    if bold is not None: run.font.bold = bold

    # Subsequent paragraphs
    for ln in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        if font_name: r.font.name = font_name
        if font_size: r.font.size = font_size

def fill_table(table, content_dict, value_col=1):
    """Fill a table where col0 is the row label and col[value_col] is to be filled."""
    for row in table.rows:
        label = row.cells[0].text.strip().replace("\n", " ").replace("|", "").strip()
        # Try various normalizations
        normalized = " ".join(label.split())
        for key, val in content_dict.items():
            key_norm = " ".join(key.split())
            if normalized.lower() == key_norm.lower():
                set_cell_text(row.cells[value_col], val)
                break
            # also handle the joined-with-pipes in source
            if normalized.replace(" ", "").lower() == key_norm.replace(" ", "").lower():
                set_cell_text(row.cells[value_col], val)
                break

def set_shape_text(shape, text):
    """Replace a text box's content, preserving first run's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        bold = first_run.font.bold
        color = first_run.font.color.rgb if first_run.font.color and first_run.font.color.type else None
    else:
        font_name = font_size = bold = color = None

    txBody = tf._txBody
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    paras = txBody.findall('a:p', nsmap)
    for p in paras[1:]:
        txBody.remove(p)
    first_p = paras[0] if paras else None
    if first_p is not None:
        for r in first_p.findall('a:r', nsmap):
            first_p.remove(r)
        for br in first_p.findall('a:br', nsmap):
            first_p.remove(br)

    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    run = p0.add_run()
    run.text = lines[0]
    if font_name: run.font.name = font_name
    if font_size: run.font.size = font_size
    if bold is not None: run.font.bold = bold

    for ln in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        if font_name: r.font.name = font_name
        if font_size: r.font.size = font_size

# ============ MAIN ============

def main():
    prs = Presentation(SRC)

    # Slide 1 — Team & Challenge Statement (5x3 table at shape[2])
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_table:
            fill_table(shape.table, SLIDE1, value_col=1)

    # Slide 2 — Project Proposal etc. (4x3 table)
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_table:
            fill_table(shape.table, SLIDE2, value_col=1)

    # Slide 3 — Technical Architecture (4x2 table)
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_table:
            fill_table(shape.table, SLIDE3, value_col=1)

    # Slide 4 — Team Profile (6 NAME shapes)
    slide4 = prs.slides[3]
    name_idx = 0
    role_idx = 0
    name_shapes = []
    for shape in slide4.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "<NAME>":
            name_shapes.append(shape)
    for i, shape in enumerate(name_shapes):
        if i < len(TEAM):
            set_shape_text(shape, TEAM[i][0])
    # Look for the role/title shapes near each NAME — they're likely the empty shapes
    # immediately before each NAME. Iterate all shapes; whenever we see <NAME>-was, use the previous empty one.
    # Simpler approach: find empty AUTO_SHAPE shapes adjacent to NAME shapes.
    all_shapes = list(slide4.shapes)
    name_indices = [i for i, s in enumerate(all_shapes) if s.has_text_frame and s.text_frame.text.strip() in [t[0] for t in TEAM] + ["<NAME>"]]
    # The pattern in the template appears to be: 4 empty shapes per name (decorative), then NAME
    # Find empty text-frame shapes RIGHT AFTER each NAME shape and put the role there
    for member_idx, name_idx in enumerate(name_indices):
        if member_idx >= len(TEAM):
            break
        # Look at the next shape after the NAME
        for offset in (1, 2, 3, -1, -2):
            target_idx = name_idx + offset
            if 0 <= target_idx < len(all_shapes):
                target = all_shapes[target_idx]
                if target.has_text_frame and target.text_frame.text.strip() == "":
                    set_shape_text(target, TEAM[member_idx][1])
                    break

    # Slide 5 — References (find the empty text frame shape that's not the title)
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "" or txt == "Other References/Points (If Any)":
                if txt == "":
                    set_shape_text(shape, SLIDE5_REFS)
                    break

    prs.save(DST)
    print(f"Wrote {DST}")

if __name__ == "__main__":
    main()
